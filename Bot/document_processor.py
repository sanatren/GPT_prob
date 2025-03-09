import os
import tempfile
import streamlit as st
from typing import List, Dict, Any, Optional
import docx
import pandas as pd
from pptx import Presentation
import pdfplumber
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.schema import Document

class DocumentProcessor:
    """Process various document types for RAG applications"""
    
    def __init__(self):
        """Initialize the document processor"""
        # Initialize the text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        
        # Initialize the embeddings model
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2"
        )
        
        # Dictionary to store vectorstores by session ID
        self.vectorstores = {}
    
    def process_file(self, uploaded_file, session_id: str) -> bool:
        """
        Process an uploaded file and add it to the vector store
        
        Args:
            uploaded_file: The uploaded file from Streamlit
            session_id: The current session ID
            
        Returns:
            bool: True if processing was successful, False otherwise
        """
        try:
            # Save the uploaded file to a temporary location
            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{uploaded_file.name.split(".")[-1]}') as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                file_path = tmp_file.name
            
            # Extract text based on file type
            file_extension = uploaded_file.name.split('.')[-1].lower()
            
            if file_extension == 'pdf':
                text = self._extract_text_from_pdf(file_path)
            elif file_extension == 'docx':
                text = self._extract_text_from_docx(file_path)
            elif file_extension in ['xlsx', 'xls']:
                text = self._extract_text_from_excel(file_path)
            elif file_extension == 'pptx':
                text = self._extract_text_from_pptx(file_path)
            elif file_extension == 'txt':
                text = self._extract_text_from_txt(file_path)
            else:
                st.error(f"Unsupported file type: {file_extension}")
                os.unlink(file_path)  # Clean up the temp file
                return False
            
            # Clean up the temp file
            os.unlink(file_path)
            
            if not text:
                st.warning("No text could be extracted from the file.")
                return False
            
            # Create metadata for the document
            metadata = {
                "source": uploaded_file.name,
                "file_type": file_extension
            }
            
            # Create a LangChain Document
            doc = Document(page_content=text, metadata=metadata)
            
            # Split the document into chunks
            chunks = self.text_splitter.split_documents([doc])
            
            # Create or update the vector store for this session
            if session_id in self.vectorstores:
                # Add to existing vectorstore
                self.vectorstores[session_id].add_documents(chunks)
            else:
                # Create new vectorstore
                self.vectorstores[session_id] = FAISS.from_documents(
                    chunks, self.embeddings
                )
            
            return True
            
        except Exception as e:
            st.error(f"Error processing file: {str(e)}")
            return False
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file"""
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text
    
    def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from a Word document"""
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from an Excel file"""
        df = pd.read_excel(file_path, sheet_name=None)
        text = ""
        
        # Process each sheet
        for sheet_name, sheet_df in df.items():
            text += f"\n\nSheet: {sheet_name}\n"
            
            # Convert the dataframe to a string representation
            text += sheet_df.to_string(index=False)
        
        return text
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from a PowerPoint presentation"""
        prs = Presentation(file_path)
        text = ""
        
        for i, slide in enumerate(prs.slides):
            text += f"\n\nSlide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text
    
    def _extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from a text file"""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            return file.read()
    
    def query_documents(self, query: str, session_id: str, k: int = 4) -> List[Dict[str, Any]]:
        """
        Query the vector store for relevant document chunks
        
        Args:
            query: The user's query
            session_id: The current session ID
            k: Number of results to return
            
        Returns:
            List of dictionaries containing content and metadata
        """
        if session_id not in self.vectorstores:
            return []
        
        vectorstore = self.vectorstores[session_id]
        results = vectorstore.similarity_search_with_score(query, k=k)
        
        # Format the results
        formatted_results = []
        for doc, score in results:
            formatted_results.append({
                "content": doc.page_content,
                "metadata": doc.metadata,
                "score": float(score)
            })
        
        return formatted_results
    
    def clear_documents(self, session_id: str) -> bool:
        """Clear all documents for a session"""
        if session_id in self.vectorstores:
            del self.vectorstores[session_id]
            return True
        return False 